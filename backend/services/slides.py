from typing import List, Tuple, Dict, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import subprocess
from utils.paths import with_outputs

# ===== 預設主題（可依需求新增） =====
PRESET_THEMES: Dict[str, Dict] = {
    "clean_light": {
        "ratio": "16:9",
        "bg_color": "#FFFFFF",
        "title_color": "#111111",
        "body_color": "#222222",
        "accent_color": "#3B82F6",
        "title_font": "PingFang TC",
        "body_font": "PingFang TC",
        "title_size": 40,
        "body_size": 22,
        "footer_text": ""
    },
    "clean_dark": {
        "ratio": "16:9",
        "bg_color": "#0B1221",
        "title_color": "#E5E7EB",
        "body_color": "#D1D5DB",
        "accent_color": "#60A5FA",
        "title_font": "PingFang TC",
        "body_font": "PingFang TC",
        "title_size": 42,
        "body_size": 24,
        "footer_text": ""
    },
    "corporate_blue": {
        "ratio": "16:9",
        "bg_color": "#F8FAFF",
        "title_color": "#0F172A",
        "body_color": "#1E293B",
        "accent_color": "#2563EB",
        "title_font": "Arial",
        "body_font": "Arial",
        "title_size": 40,
        "body_size": 22,
        "footer_text": "Confidential"
    }
}

def list_theme_presets() -> List[Dict]:
    return [{"name": k, **v} for k, v in PRESET_THEMES.items()]

def _hex_to_rgb(hex_str: str) -> RGBColor:
    hs = hex_str.lstrip('#')
    return RGBColor(int(hs[0:2], 16), int(hs[2:4], 16), int(hs[4:6], 16))

def _apply_ratio(prs: Presentation, ratio: str):
    if ratio == "4:3":
        prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    else:  # 16:9 預設
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

def _apply_bg(slide, theme: Dict):
    bg_img = theme.get("bg_image_path")
    if bg_img:
        p = Path(bg_img)
        if p.exists():
            pic = slide.shapes.add_picture(str(p), 0, 0, width=slide.part.slide_width, height=slide.part.slide_height)
            pic.zorder = 0
            return
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(theme.get("bg_color", "#FFFFFF"))

def _maybe_add_logo(slide, theme: Dict):
    logo = theme.get("logo_path")
    if logo:
        p = Path(logo)
        if p.exists():
            w = Inches(1.1)
            left = slide.part.slide_width - w - Inches(0.3)
            top = Inches(0.2)
            slide.shapes.add_picture(str(p), left, top, width=w, height=None)

def _maybe_footer(slide, theme: Dict):
    text = theme.get("footer_text")
    if text:
        box = slide.shapes.add_textbox(Inches(0.5), slide.part.slide_height - Inches(0.5), slide.part.slide_width - Inches(1.0), Inches(0.3))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        font = run.font
        font.size = Pt(10)
        font.name = theme.get("body_font", "Arial")
        font.color.rgb = _hex_to_rgb(theme.get("body_color", "#222222"))
        p.alignment = PP_ALIGN.RIGHT

def _style_title_placeholder(slide, theme: Dict):
    if not slide.shapes.title:
        return
    title_shape = slide.shapes.title
    title_tf = title_shape.text_frame
    for p in title_tf.paragraphs:
        for r in p.runs:
            f = r.font
            f.name = theme.get("title_font", "Arial")
            f.size = Pt(theme.get("title_size", 40))
            f.color.rgb = _hex_to_rgb(theme.get("title_color", "#111111"))
        p.alignment = PP_ALIGN.LEFT

def _style_body_placeholder(slide, theme: Dict):
    body = None
    for shp in slide.placeholders:
        if shp.placeholder_format.idx == 1:
            body = shp
            break
    if not body:
        return
    tf = body.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        if not p.runs:
            r = p.add_run(); r.text = p.text
        for r in p.runs:
            f = r.font
            f.name = theme.get("body_font", "Arial")
            f.size = Pt(theme.get("body_size", 22))
            f.color.rgb = _hex_to_rgb(theme.get("body_color", "#222222"))
        p.level = 0

def _merge_theme(preset_name: Optional[str], overrides: Optional[Dict]) -> Dict:
    base = PRESET_THEMES.get(preset_name or "clean_light", PRESET_THEMES["clean_light"]).copy()
    if overrides:
        overrides = {k: v for k, v in overrides.items() if k != "preset_name" and v is not None}
        base.update(overrides)
    return base

def _to_pdf_with_libreoffice(pptx_path: Path) -> Path:
    outdir = with_outputs("").parent
    cmd = ["soffice", "--headless", "--convert-to", "pdf", str(pptx_path), "--outdir", str(outdir)]
    try:
        import subprocess as _sp
        _sp.run(cmd, check=True, stdout=_sp.PIPE, stderr=_sp.PIPE)
        pdf_path = outdir / (pptx_path.stem + ".pdf")
        return pdf_path
    except Exception as e:
        raise RuntimeError("LibreOffice 未安裝或轉檔失敗") from e

def make_slides(title: str, points: List[str], preset_name: Optional[str] = None, theme_overrides: Optional[Dict] = None) -> Tuple[Path, Path]:
    theme = _merge_theme(preset_name, theme_overrides)

    prs = Presentation()
    _apply_ratio(prs, theme.get("ratio", "16:9"))

    # 標題頁
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = "Auto-generated briefing"
    _apply_bg(slide, theme)
    _style_title_placeholder(slide, theme)
    _style_body_placeholder(slide, theme)
    _maybe_add_logo(slide, theme)
    _maybe_footer(slide, theme)

    # 內容頁（每頁最多 6 點）
    chunk = 6
    for i in range(0, len(points), chunk):
        layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Key Points ({i+1}-{min(i+chunk, len(points))})"
        body = slide.placeholders[1]
        body.text = "\n".join([f"• {p}" for p in points[i:i+chunk]])
        _apply_bg(slide, theme)
        _style_title_placeholder(slide, theme)
        _style_body_placeholder(slide, theme)
        _maybe_add_logo(slide, theme)
        _maybe_footer(slide, theme)

    pptx_path = with_outputs("auto_slides.pptx")
    prs.save(pptx_path)
    pdf_path = _to_pdf_with_libreoffice(pptx_path)
    return pptx_path, pdf_path
