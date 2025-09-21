import os, shlex
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from .utils import run_cmd

def _add_title_slide(prs, title):
    layout = prs.slide_layouts[0]  # title slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Auto‑generated summary"
    return slide

def _add_bullets_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title[:80]
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for b in bullets[:8]:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
    return slide

def make_slides(summary: Dict[str, Any], out_pptx: str, out_pdf: str):
    prs = Presentation()
    title = "文件重點 / Key Points" if summary.get("language","auto") in ("zh","zh-hant","zh-hans","auto") else "Key Points"
    _add_title_slide(prs, title)
    for pt in summary.get("points", []):
        t = pt.get("title", "Point")
        bullets = pt.get("bullets", [])
        _add_bullets_slide(prs, t, bullets)
    prs.save(out_pptx)

    # Convert PPTX → PDF via LibreOffice
    cmd = f"soffice --headless --convert-to pdf --outdir {shlex.quote(os.path.dirname(out_pdf))} {shlex.quote(out_pptx)}"
    run_cmd(cmd)
    # LO names output same base
    # if file at desired name missing, try to find any pdf in outdir
    if not os.path.exists(out_pdf):
        base = os.path.splitext(os.path.basename(out_pptx))[0] + ".pdf"
        alt = os.path.join(os.path.dirname(out_pdf), base)
        if os.path.exists(alt):
            os.rename(alt, out_pdf)
